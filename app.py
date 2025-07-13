import streamlit as st
import base64
import os
from pptx import Presentation
from pptx.util import Pt
from transformers import pipeline

# Format settings
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

# Load Hugging Face pipeline (gpt2 is free, may use other open-source models)
@st.cache_resource
def get_generator():
    return pipeline("text-generation", model="gpt2")

generator = get_generator()

# Generate slide titles (simulate outline generation)
def generate_slide_titles(topic):
    prompt = f"List 5 informative and concise slide titles for a PowerPoint presentation about: '{topic}'."
    result = generator(prompt, max_length=80, num_return_sequences=1)[0]["generated_text"]
    # Extract titles from text, splitting by common delimiters
    lines = result.split('\n')
    titles = [line.strip("-•1234567890. ").strip() for line in lines if line.strip()]
    # Take only first 5 unique, non-empty titles
    titles = [t for t in titles if t][:5]
    return titles

# Generate slide content
def generate_slide_content(title):
    prompt = f"Write a concise PowerPoint slide paragraph (4–5 sentences) about: '{title}'."
    result = generator(prompt, max_length=120, num_return_sequences=1)[0]["generated_text"]
    # Clean up output: take only the new sentences after the prompt
    paragraph = result.replace(prompt, '').strip()
    return paragraph

# Create PPT
def create_ppt(topic, titles, contents):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Content Slides
    for title, content in zip(titles, contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1]
        body.text = content

        # Font customization
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    p.font.size = SLIDE_FONT_SIZE

    os.makedirs("generated_ppt", exist_ok=True)
    path = f"generated_ppt/{topic}_presentation.pptx"
    prs.save(path)
    return path

# Download link
def get_download_link(path):
    with open(path, "rb") as f:
        data = f.read()
    b64 = base64.b64encode(data).decode()
    filename = os.path.basename(path)
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">📥 Download the PowerPoint</a>'

# Streamlit App
def main():
    st.title("... Free Text-to-PPT Generator (Open-Source AI)")

    topic = st.text_input("Enter your presentation topic:")
    if st.button("Generate Presentation") and topic:
        st.info("Generating slides with open-source AI...")

        titles = generate_slide_titles(topic)
        contents = [generate_slide_content(title) for title in titles]
        path = create_ppt(topic, titles, contents)

        st.success("✅ Presentation created!")
        st.markdown(get_download_link(path), unsafe_allow_html=True)

if __name__ == "__main__":
    main()