import os
import time
from pptx import Presentation
from pptx.util import Pt
from transformers import pipeline

# Constants for font sizes
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

# Load Hugging Face pipeline (GPT-2)
def get_generator():
    return pipeline("text-generation", model="gpt2")

generator = get_generator()

def generate_slide_titles(topic):
    prompt = f"List 5 concise and professional slide titles for a PowerPoint presentation about: '{topic}'."
    result = generator(prompt, max_length=80, num_return_sequences=1)[0]["generated_text"]
    # Extract titles from text, splitting by common delimiters
    lines = result.split('\n')
    titles = [line.strip("-•1234567890. ").strip() for line in lines if line.strip()]
    titles = [t for t in titles if t][:5]
    return titles

def generate_slide_content(title):
    prompt = f"Write a clear PowerPoint slide paragraph (4–5 sentences) explaining: '{title}'."
    result = generator(prompt, max_length=120, num_return_sequences=1)[0]["generated_text"]
    paragraph = result.replace(prompt, '').strip()
    return paragraph

def create_presentation(topic, slide_titles, slide_contents):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Content Slides
    for title, content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

        # Font styling
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    p.font.size = SLIDE_FONT_SIZE

    os.makedirs("generated_ppt", exist_ok=True)
    filepath = f"generated_ppt/{topic}_presentation.pptx"
    prs.save(filepath)
    return filepath

def main():
    topic = input("Enter your presentation topic: ").strip()
    if not topic:
        print("❌ Please enter a topic.")
        return

    print(f"\n🔄 Generating slide titles for: {topic}")
    start = time.time()

    titles = generate_slide_titles(topic)
    print("✅ Slide Titles:")
    for i, title in enumerate(titles, 1):
        print(f"{i}. {title}")

    print("\n✍️ Generating slide content...")
    contents = [generate_slide_content(title) for title in titles]

    output_path = create_presentation(topic, titles, contents)

    end = time.time()
    print(f"\n✅ Presentation saved to: {output_path}")
    print(f"⏱️ Time taken: {round(end - start, 2)} seconds")

if __name__ == "__main__":
    main()